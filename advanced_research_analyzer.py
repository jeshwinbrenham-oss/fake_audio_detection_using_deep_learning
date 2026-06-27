import streamlit as st
import tempfile
import os
import re
import zipfile
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import docx
import docx2txt

# ---------------- SAFE IMPORTS ---------------- #
try:
    import pdfplumber
    PDF_OK = True
except:
    PDF_OK = False

try:
    from pptx import Presentation
    PPT_OK = True
except:
    PPT_OK = False

from transformers import pipeline
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity

# ---------------- PAGE CONFIG ---------------- #
st.set_page_config(
    page_title="Ultimate Research Analyzer",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ---------------- SIDEBAR ---------------- #
st.sidebar.title("⚙️ Analyzer Settings")

ai_threshold = st.sidebar.slider(
    "AI Detection Threshold",
    0.0,
    1.0,
    0.60,
    0.01
)

similarity_threshold = st.sidebar.slider(
    "Similarity Threshold",
    0.50,
    1.00,
    0.85,
    0.01
)

# ---------------- TITLE ---------------- #
st.title("🧠 Ultimate AI Journal Research Analyzer")
st.markdown("### Professional Journal-Level Screening System")

# ---------------- LOAD MODELS ---------------- #
@st.cache_resource
def load_models():
    ai_model = pipeline(
        "text-classification",
        model="roberta-base-openai-detector",
        return_all_scores=True
    )

    embed_model = SentenceTransformer("all-MiniLM-L6-v2")

    return ai_model, embed_model

ai_model, embed_model = load_models()

# ---------------- DOCX READING ---------------- #
def read_docx(path):
    try:
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except:
        return None


def fallback_docx(path):
    try:
        text = docx2txt.process(path)
        if text.strip():
            return text
    except:
        pass

    try:
        with zipfile.ZipFile(path) as z:
            xml = z.read("word/document.xml").decode("utf-8")
            return re.sub("<.*?>", " ", xml)
    except:
        return ""

# ---------------- FILE READER ---------------- #
def extract_text(file):
    text = ""

    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file.read())
        path = tmp.name

    name = file.name.lower()

    try:
        # DOCX
        if name.endswith(".docx"):
            text = read_docx(path)

            if text is None:
                st.warning("⚠️ DOCX damaged → recovery mode enabled")
                text = fallback_docx(path)

        # PDF
        elif name.endswith(".pdf"):
            if PDF_OK:
                with pdfplumber.open(path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
            else:
                st.error("pdfplumber not installed")

        # PPTX
        elif name.endswith(".pptx"):
            if PPT_OK:
                prs = Presentation(path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            else:
                st.error("python-pptx not installed")

        # TXT
        elif name.endswith(".txt"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

    except Exception as e:
        st.error(f"File reading error: {e}")

    os.remove(path)

    text = text.replace("\n", " ")

    sentences = re.split(r'[.!?]', text)

    sentences = [
        s.strip()
        for s in sentences
        if len(s.strip()) > 20
    ]

    return text, sentences

# ---------------- AI DETECTION ---------------- #
def ai_score(sentence):
    try:
        result = ai_model(sentence[:512])[0]
        top = max(result, key=lambda x: x['score'])
        return top['score']
    except:
        return 0

# ---------------- STYLE ANALYSIS ---------------- #
def style_score(text):
    words = text.split()

    if len(words) < 5:
        return 0

    return np.std([len(w) for w in words])

# ---------------- ANALYSIS ---------------- #
def analyze_sentences(sentences):
    results = []

    for s in sentences:
        ai = ai_score(s)
        style = style_score(s)

        final = (ai * 0.7) + (1 / (style + 1)) * 0.3

        results.append({
            "sentence": s,
            "ai_score": round(ai, 4),
            "style_score": round(style, 4),
            "final_score": round(final, 4)
        })

    return results

# ---------------- SIMILARITY ---------------- #
def semantic_similarity(sentences):
    if len(sentences) < 2:
        return []

    embeddings = embed_model.encode(sentences)
    sim = cosine_similarity(embeddings)

    pairs = []

    for i in range(len(sentences)):
        for j in range(i + 1, len(sentences)):
            if sim[i][j] > similarity_threshold:
                pairs.append({
                    "sentence1": sentences[i],
                    "sentence2": sentences[j],
                    "score": round(float(sim[i][j]), 4)
                })

    return pairs

# ---------------- STRUCTURE ---------------- #
def structure_analysis(text):
    sections = [
        "abstract",
        "introduction",
        "method",
        "methodology",
        "result",
        "results",
        "discussion",
        "conclusion",
        "references"
    ]

    found = []

    low = text.lower()

    for sec in sections:
        if sec in low:
            found.append(sec)

    return found

# ---------------- STYLE ISSUES ---------------- #
def style_issues(sentences):
    lengths = [len(s.split()) for s in sentences]

    if len(lengths) < 5:
        return []

    avg = np.mean(lengths)
    std = np.std(lengths)

    issues = []

    for i, l in enumerate(lengths):
        if abs(l - avg) > 2 * std:
            issues.append(sentences[i])

    return issues

# ---------------- REVIEWER AI ---------------- #
def reviewer_system(ai_percent, sim_count, structure_count, style_count):

    reviewers = []

    # Reviewer 1
    if ai_percent > 40:
        r1 = "Writing appears heavily AI-generated."
        s1 = 40
    elif ai_percent > 20:
        r1 = "Moderate AI-like writing detected."
        s1 = 65
    else:
        r1 = "Writing quality appears mostly human-authored."
        s1 = 90

    reviewers.append(("Reviewer A - AI Quality", r1, s1))

    # Reviewer 2
    if sim_count > 5:
        r2 = "High semantic duplication detected."
        s2 = 45
    elif sim_count > 2:
        r2 = "Some duplicated concepts detected."
        s2 = 70
    else:
        r2 = "Originality level is acceptable."
        s2 = 92

    reviewers.append(("Reviewer B - Originality", r2, s2))

    # Reviewer 3
    if structure_count < 4:
        r3 = "Research structure is incomplete."
        s3 = 50
    else:
        r3 = "Research structure is professionally organized."
        s3 = 90

    reviewers.append(("Reviewer C - Structure", r3, s3))

    # Reviewer 4
    if style_count > 10:
        r4 = "Academic writing consistency is weak."
        s4 = 55
    else:
        r4 = "Writing style consistency is acceptable."
        s4 = 88

    reviewers.append(("Reviewer D - Academic Tone", r4, s4))

    return reviewers

# ---------------- PUBLICATION DECISION ---------------- #
def editorial_decision(ai_percent, sim_count, structure_count, style_count):

    if ai_percent > 60 or sim_count > 10 or structure_count < 3:
        return "❌ REJECTED"

    elif ai_percent > 35 or sim_count > 5 or style_count > 10:
        return "⚠️ MAJOR REVISION"

    elif ai_percent > 15 or style_count > 5:
        return "⚠️ MINOR REVISION"

    else:
        return "✅ ACCEPTED"

# ---------------- PROBABILITY ---------------- #
def acceptance_probability(score):
    return max(min(score, 100), 0)

# ---------------- FILE UPLOADER ---------------- #
file = st.file_uploader(
    "📄 Upload Research Paper",
    type=["docx", "pdf", "pptx", "txt"]
)

# ---------------- MAIN EXECUTION ---------------- #
if file:

    with st.spinner("Analyzing research paper..."):

        full_text, sentences = extract_text(file)

        if not sentences:
            st.error("❌ No readable text found")
            st.stop()

        st.success(f"✅ Extracted {len(sentences)} sentences")

        # AI Analysis
        results = analyze_sentences(sentences)

        ai_count = sum(
            1
            for r in results
            if r['final_score'] > ai_threshold
        )

        ai_percent = (ai_count / len(sentences)) * 100

        # Similarity
        sim_pairs = semantic_similarity(sentences)

        # Structure
        found_sections = structure_analysis(full_text)
        structure_score = len(found_sections)

        # Style
        issues = style_issues(sentences)

        # Final Score
        score = 100

        score -= ai_percent * 0.5
        score -= len(sim_pairs) * 2
        score -= len(issues) * 1.5
        score += structure_score * 3

        score = round(max(min(score, 100), 0), 2)

        # Decision
        decision = editorial_decision(
            ai_percent,
            len(sim_pairs),
            structure_score,
            len(issues)
        )

        # Probability
        publication_prob = acceptance_probability(score)

        # Reviewers
        reviewers = reviewer_system(
            ai_percent,
            len(sim_pairs),
            structure_score,
            len(issues)
        )

    # ---------------- DASHBOARD ---------------- #
    st.subheader("📊 Editorial Metrics")

    c1, c2, c3, c4 = st.columns(4)

    c1.metric("AI Risk", f"{round(ai_percent,2)}%")
    c2.metric("Similarity", len(sim_pairs))
    c3.metric("Style Issues", len(issues))
    c4.metric("Structure", f"{structure_score}/9")

    st.progress(score / 100)

    # ---------------- DECISION ---------------- #
    st.subheader("📢 Editorial Decision")

    if "REJECTED" in decision:
        st.error(decision)
    elif "REVISION" in decision:
        st.warning(decision)
    else:
        st.success(decision)

    st.write(f"### 🎯 Reviewer Score: {score}/100")
    st.write(f"### 📈 Publication Probability: {publication_prob}%")

    # ---------------- CHART ---------------- #
    st.subheader("📉 Risk Analysis")

    chart_df = pd.DataFrame({
        "Metric": [
            "AI Risk",
            "Similarity",
            "Style Issues",
            "Structure"
        ],
        "Value": [
            ai_percent,
            len(sim_pairs),
            len(issues),
            structure_score
        ]
    })

    fig, ax = plt.subplots(figsize=(8,4))
    ax.bar(chart_df["Metric"], chart_df["Value"])
    st.pyplot(fig)

    # ---------------- REVIEWERS ---------------- #
    st.subheader("🧾 AI Reviewer Panel")

    for name, comment, rating in reviewers:
        with st.expander(f"{name} ({rating}/100)"):
            st.write(comment)

    # ---------------- SECTIONS ---------------- #
    st.subheader("📚 Detected Research Sections")

    for sec in found_sections:
        st.success(sec.upper())

    # ---------------- HIGH AI ---------------- #
    st.subheader("🚨 High AI Risk Sentences")

    high_ai = [
        r for r in results
        if r['final_score'] > 0.80
    ]

    if high_ai:
        for r in high_ai[:20]:
            st.error(r['sentence'])
            st.caption(f"AI Score: {r['final_score']}")
    else:
        st.success("No dangerous AI-like sentences detected")

    # ---------------- SIMILARITY ---------------- #
    st.subheader("🔁 Semantic Similarity Detection")

    if sim_pairs:
        for pair in sim_pairs[:10]:
            st.warning(f"Similarity: {pair['score'] * 100:.2f}%")
            st.write("Sentence 1:")
            st.info(pair['sentence1'])
            st.write("Sentence 2:")
            st.info(pair['sentence2'])
            st.divider()
    else:
        st.success("No major similarity detected")

    # ---------------- STYLE ISSUES ---------------- #
    st.subheader("⚠️ Writing Style Issues")

    if issues:
        for item in issues[:10]:
            st.warning(item)
    else:
        st.success("Academic style consistency looks good")

    # ---------------- FINAL COMMENTS ---------------- #
    st.subheader("📝 Final Editorial Comments")

    comments = []

    if ai_percent > 30:
        comments.append("Large portions may appear AI-generated.")

    if len(sim_pairs) > 3:
        comments.append("Potential semantic duplication detected.")

    if structure_score < 5:
        comments.append("Research paper structure is incomplete.")

    if len(issues) > 5:
        comments.append("Academic writing consistency should be improved.")

    if decision == "✅ ACCEPTED":
        comments.append("Paper appears professionally prepared for submission.")

    if comments:
        for c in comments:
            st.write(f"• {c}")

    # ---------------- DOWNLOAD ---------------- #
    st.subheader("📥 Export Report")

    report = f'''
ULTIMATE RESEARCH ANALYZER REPORT

AI Risk: {ai_percent}%
Similarity Count: {len(sim_pairs)}
Style Issues: {len(issues)}
Structure Score: {structure_score}/9

Editorial Decision: {decision}
Publication Probability: {publication_prob}%
Reviewer Score: {score}/100
'''

    st.download_button(
        label="📄 Download Analysis Report",
        data=report,
        file_name="research_report.txt",
        mime="text/plain"
    )

else:
    st.info("📄 Upload a research paper to start analysis")